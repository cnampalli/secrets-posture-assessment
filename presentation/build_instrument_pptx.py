#!/usr/bin/env python3
"""
Build a printable, editable 16:9 PowerPoint of the "Instrument — One Posture" slide.

Mirrors presentation/instrument-one-slide.html in a deck-native form:
dark obsidian background, gold/emerald accents, convergence concept, a
four-cell business-value ledger, and a regulatory footer.

Run:  .venv-pptx/bin/python build_instrument_pptx.py
Out:  presentation/instrument-one-slide.pptx
"""
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import instrument_stats  # noqa: E402

# every headline count is machine-derived from the live data — never hardcoded
STATS = instrument_stats.compute_stats()

# ---- palette (matches the HTML) ----
BG      = RGBColor(0x07, 0x0A, 0x10)
BG2     = RGBColor(0x12, 0x1A, 0x28)
INK     = RGBColor(0xF4, 0xF1, 0xEA)
INK2    = RGBColor(0xC8, 0xCD, 0xD6)
MUTED   = RGBColor(0x8A, 0x93, 0xA4)
HINT    = RGBColor(0x5C, 0x65, 0x77)
EMERALD = RGBColor(0x1F, 0xB2, 0x87)
CYAN    = RGBColor(0x19, 0xD7, 0xC6)
GOLD    = RGBColor(0xE8, 0xB0, 0x4B)
INDIGO  = RGBColor(0x61, 0x72, 0xFF)
CARD    = RGBColor(0x0E, 0x14, 0x1D)
LINE    = RGBColor(0x2A, 0x32, 0x40)

DISPLAY = "Georgia"          # serif display (Instrument Serif not guaranteed on host)
SANS    = "Hanken Grotesk"   # falls back gracefully if absent
MONO    = "Consolas"

EMU_IN = 914400


def solid(shape, color, line_color=None, line_w=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w or 0.75)
    shape.shadow.inherit = False
    return shape


def grad_fill(shape, c1, c2, angle=90):
    """Two-stop linear gradient via direct XML (python-pptx has no high-level API)."""
    spPr = shape.fill._xPr  # noqa: SLF001
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        e = spPr.find(qn(tag))
        if e is not None:
            spPr.remove(e)
    grad = spPr.makeelement(qn("a:gradFill"), {})
    lst = grad.makeelement(qn("a:gsLst"), {})
    for pos, col in ((0, c1), (100000, c2)):
        gs = grad.makeelement(qn("a:gs"), {"pos": str(pos)})
        clr = grad.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % (col[0], col[1], col[2])})
        gs.append(clr)
        lst.append(gs)
    grad.append(lst)
    lin = grad.makeelement(qn("a:lin"), {"ang": str(int(angle * 60000)), "scaled": "1"})
    grad.append(lin)
    # insert gradient before a:ln if present
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        ln.addprevious(grad)
    else:
        spPr.append(grad)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb, tf


def run(p, text, size, color, font=SANS, bold=False, italic=False, spacing=None):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.name = font
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if spacing is not None:  # letter spacing in points
        rPr = r._r.get_or_add_rPr()
        rPr.set("spc", str(int(spacing * 100)))
    return r


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = 13.333, 7.5
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ---- background gradient ----
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    grad_fill(bg, (0x0B, 0x12, 0x1F), (0x06, 0x08, 0x0E), angle=125)

    # soft aurora accents (large translucent-ish circles using deep tints)
    for x, y, d, col in [(-2.2, -2.4, 7.2, (0x12, 0x33, 0x2A)),
                         (9.2, -2.0, 6.6, (0x18, 0x1E, 0x3E)),
                         (8.4, 4.6, 6.8, (0x33, 0x2A, 0x16))]:
        o = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
        solid(o, RGBColor(*col))

    # top hairline
    hl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(0.92), Inches(SW - 1.8), Pt(1))
    grad_fill(hl, (0xE8, 0xB0, 0x4B), (0x07, 0x0A, 0x10), angle=0)

    # ---- header ----
    gem = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(0.9), Inches(0.5), Inches(0.16), Inches(0.16))
    grad_fill(gem, (0xE8, 0xB0, 0x4B), (0x1F, 0xB2, 0x87), angle=135)
    _, tf = textbox(slide, 1.18, 0.46, 6, 0.3, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    run(p, "THE ", 12, INK2, MONO, spacing=3)
    run(p, "INSTRUMENT", 12, INK, MONO, bold=True, spacing=3)
    _, tf = textbox(slide, SW - 6.4, 0.46, 5.5, 0.3, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    run(p, "IDENTITY-SECURITY POSTURE · VENDOR SELECTION", 8.5, MUTED, MONO, spacing=2)

    # ---- eyebrow ----
    _, tf = textbox(slide, 0.9, 1.35, 7, 0.3)
    p = tf.paragraphs[0]
    run(p, "ONE INSTRUMENT  ·  ONE POSTURE", 10, GOLD, MONO, bold=True, spacing=2.5)

    # ---- headline ----
    _, tf = textbox(slide, 0.86, 1.7, 7.2, 2.2)
    p = tf.paragraphs[0]; p.line_spacing = 0.98
    run(p, "One verdict you can ", 40, INK, DISPLAY)
    run(p, "defend", 40, GOLD, DISPLAY, italic=True)
    p2 = tf.add_paragraph(); p2.line_spacing = 0.98
    run(p2, "across the whole identity estate.", 40, INK, DISPLAY)

    # ---- lede ----
    _, tf = textbox(slide, 0.9, 3.95, 6.7, 1.5)
    p = tf.paragraphs[0]; p.line_spacing = 1.3
    run(p, "An ", 14.5, INK2, SANS)
    run(p, "offline, evidence-gated", 14.5, INK, SANS, bold=True)
    run(p, " engine that scores Secrets, Privileged Access & Identity "
            "Governance as a single posture — and proves every verdict, with a cited "
            "source, to your ", 14.5, INK2, SANS)
    run(p, "board and your regulator.", 14.5, INK, SANS, bold=True)

    # ---- pills ----
    pills = [("Self-hosted", EMERALD), ("Evidence-gated", GOLD),
             ("Board-defensible", CYAN), ("2026-forward taxonomy", INDIGO)]
    px = 0.9
    for label, col in pills:
        w = 0.42 + 0.092 * len(label)
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(px), Inches(5.35), Inches(w), Inches(0.42))
        pill.adjustments[0] = 0.5
        solid(pill, CARD, line_color=LINE, line_w=1)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px + 0.16), Inches(5.51), Inches(0.1), Inches(0.1))
        solid(dot, col)
        _, tf = textbox(slide, px + 0.32, 5.36, w - 0.34, 0.4, MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        run(p, label.upper(), 8.5, INK2, MONO, spacing=1)
        px += w + 0.18

    # ---- convergence orb (right) ----
    cx, cy, R = 10.55, 3.0, 1.62
    ring1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - R), Inches(cy - R), Inches(2 * R), Inches(2 * R))
    ring1.fill.background(); ring1.line.color.rgb = EMERALD; ring1.line.width = Pt(1.5); ring1.shadow.inherit = False
    R2 = R * 0.72
    ring2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - R2), Inches(cy - R2), Inches(2 * R2), Inches(2 * R2))
    ring2.fill.background(); ring2.line.color.rgb = INDIGO; ring2.line.width = Pt(1); ring2.shadow.inherit = False
    # core
    Rc = R * 0.46
    core = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - Rc), Inches(cy - Rc), Inches(2 * Rc), Inches(2 * Rc))
    grad_fill(core, (0x16, 0x20, 0x2E), (0x09, 0x0D, 0x15), angle=120)
    core.line.color.rgb = LINE; core.line.width = Pt(1)
    _, tf = textbox(slide, cx - Rc, cy - 0.42, 2 * Rc, 0.85, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run(p, f"{STATS['domains']} → 1", 30, INK, DISPLAY)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    run(p2, "DOMAINS · ONE POSTURE", 7, GOLD, MONO, spacing=1.5)

    # domain nodes around the orb
    nodes = [("Secrets", f"{STATS['uc_secrets']} use cases", EMERALD, cx - 0.75, cy - R - 0.42),
             ("Priv. Access", f"{STATS['uc_pam']} use cases", INDIGO, cx - R - 0.55, cy + R - 0.35),
             ("Governance", f"{STATS['uc_iga']} use cases", GOLD, cx + R - 0.95, cy + R - 0.35)]
    for nm, ct, col, nx, ny in nodes:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(nx), Inches(ny), Inches(1.55), Inches(0.62))
        box.adjustments[0] = 0.22
        solid(box, RGBColor(0x0C, 0x12, 0x1B), line_color=col, line_w=1.25)
        _, tf = textbox(slide, nx, ny + 0.06, 1.55, 0.52, MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run(p, nm, 11, INK, SANS, bold=True)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        run(p2, ct.upper(), 7, MUTED, MONO, spacing=1)

    # ---- business-value ledger (4 cells) ----
    ly, lh, lx0, gap = 5.95, 1.05, 0.9, 0.12
    lw = (SW - 1.8 - 3 * gap) / 4
    cells = [(str(STATS["uc_total"]), f"USE CASES · {STATS['archetypes']} ARCHETYPES", "Don't mis-buy",
              "Stop paying for tools that don't close your gaps.", EMERALD),
             (str(STATS["identity_total"]), "IDENTITY CLASSES", "Don't over-build",
              "See where in-house controls replace licence spend.", CYAN),
             (str(STATS["frameworks_control_level"]), "FRAMEWORKS · CONTROL-LEVEL", "Don't fail audit",
              "Cited control mappings cut audit friction and fines.", GOLD),
             ("1", "RE-RUNNABLE SOURCE OF TRUTH", "Don't re-assess",
              "One instrument replaces three RFP cycles.", INDIGO)]
    for i, (n, u, h, body, col) in enumerate(cells):
        x = lx0 + i * (lw + gap)
        cell = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(ly), Inches(lw), Inches(lh))
        cell.adjustments[0] = 0.05
        solid(cell, CARD, line_color=LINE, line_w=1)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.12), Inches(ly + 0.12), Inches(0.55), Pt(2.4))
        solid(bar, col)
        _, tf = textbox(slide, x + 0.18, ly + 0.2, lw - 0.32, lh - 0.3)
        p = tf.paragraphs[0]
        run(p, n + "  ", 22, INK, DISPLAY)
        run(p, u, 6.5, MUTED, MONO, spacing=1)
        p2 = tf.add_paragraph(); p2.space_before = Pt(4)
        run(p2, h, 11.5, INK, SANS, bold=True)
        p3 = tf.add_paragraph(); p3.space_before = Pt(2); p3.line_spacing = 1.15
        run(p3, body, 8.5, MUTED, SANS)

    # ---- footer ----
    _, tf = textbox(slide, 0.9, 7.12, 5, 0.3, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    run(p, "HUMAN · MACHINE · AGENTIC IDENTITIES", 8, HINT, MONO, spacing=2)
    _, tf = textbox(slide, SW - 8.4, 7.12, 7.5, 0.3, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    run(p, "ESSENTIAL 8 · APRA CPS 234/230 · ASD ISM · CISA ZTMM · MITRE ATT&CK",
        8, MUTED, MONO, spacing=1.5)

    out = __file__.rsplit("/", 1)[0] + "/instrument-one-slide.pptx"
    prs.save(out)
    print("wrote", out)


if __name__ == "__main__":
    build()
